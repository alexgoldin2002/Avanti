#!/usr/bin/env python3
"""Generate Avanti investor deck (.pptx) — Uber-style narrative, Avanti brand."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "Avanti-Investor-Deck.pptx"
SCREENSHOTS = ROOT / "screenshots"

# Brand tokens from app/globals.css
BRAND = {
    "forest": RGBColor(0x2A, 0x45, 0x38),
    "forest_deep": RGBColor(0x1F, 0x33, 0x29),
    "cream": RGBColor(0xFC, 0xF9, 0xF4),
    "forest_pale": RGBColor(0xE8, 0xF5, 0xEE),
    "ivory": RGBColor(0xF5, 0xF2, 0xEC),
    "muted": RGBColor(0x6A, 0x6A, 0x6A),
    "border": RGBColor(0xE0, 0xE0, 0xD8),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

SERIF = "Cormorant Garamond"
SANS = "Inter"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.75)
CONTENT_W = Inches(11.8)
FOOTER_Y = Inches(7.05)


def c(key):
    return BRAND[key]


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_run(paragraph, text, *, size=16, color=None, bold=False, italic=False, font=SERIF):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color or c("forest_deep")
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return run


def add_textbox(slide, left, top, width, height, text, **kwargs):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = kwargs.pop("align", PP_ALIGN.LEFT)
    add_run(p, text, **kwargs)
    return tf


def add_eyebrow(slide, text, top=Inches(0.42), color=None, on_dark=False):
    add_textbox(
        slide,
        MARGIN_L,
        top,
        CONTENT_W,
        Inches(0.35),
        text.upper(),
        size=10,
        color=color or (c("cream") if on_dark else c("muted")),
        font=SANS,
        align=PP_ALIGN.LEFT,
    )


def add_divider(slide, top, width=CONTENT_W, left=MARGIN_L, color=None):
    line = slide.shapes.add_shape(1, left, top, width, Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = color or c("forest")
    line.line.fill.background()


def add_footer(slide, on_dark=False):
    add_textbox(
        slide,
        MARGIN_L,
        FOOTER_Y,
        CONTENT_W,
        Inches(0.3),
        "AVANTI  ·  Confidential",
        size=9,
        color=c("cream") if on_dark else c("muted"),
        font=SANS,
        align=PP_ALIGN.RIGHT,
    )


def add_bullets(slide, left, top, width, height, items, *, size=15, color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        add_run(p, item, size=size, color=color or c("forest_deep"), font=SANS)
    return tf


def add_callout_box(slide, left, top, width, height, title, body, title_size=14, body_size=13):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = c("forest_pale")
    shape.line.color.rgb = c("border")
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p0 = tf.paragraphs[0]
    add_run(p0, title, size=title_size, bold=True, font=SANS)
    p1 = tf.add_paragraph()
    add_run(p1, body, size=body_size, color=c("muted"), font=SANS)


def slide_header(slide, eyebrow, title, subtitle=None):
    set_slide_bg(slide, c("cream"))
    add_eyebrow(slide, eyebrow)
    add_textbox(slide, MARGIN_L, Inches(0.72), CONTENT_W, Inches(0.65), title, size=34, color=c("forest_deep"))
    add_divider(slide, Inches(1.35))
    top = Inches(1.55)
    if subtitle:
        add_textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.5), subtitle, size=16, italic=True, color=c("muted"), font=SANS)
        top = Inches(2.05)
    add_footer(slide)
    return top


def headline_slide(prs, eyebrow, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = slide_header(slide, eyebrow, title, subtitle)
    add_bullets(slide, MARGIN_L, top, CONTENT_W, Inches(4.8), bullets)


def quote_slide(prs, eyebrow, title, quote, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = slide_header(slide, eyebrow, title)
    add_callout_box(slide, MARGIN_L, top, CONTENT_W, Inches(0.95), quote, "", title_size=16, body_size=14)
    add_bullets(slide, MARGIN_L, top + Inches(1.15), CONTENT_W, Inches(3.8), bullets)


def differentiators_slide(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = slide_header(slide, "Why us?", "Key differentiators")
    y = top
    for num, title, desc in items:
        badge = slide.shapes.add_shape(1, MARGIN_L, y, Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = c("forest_deep")
        badge.line.fill.background()
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        add_run(bp, num, size=12, color=c("cream"), font=SERIF)
        add_textbox(slide, MARGIN_L + Inches(0.6), y - Inches(0.02), Inches(4.5), Inches(0.35), title, size=17)
        add_textbox(
            slide,
            MARGIN_L + Inches(0.6),
            y + Inches(0.32),
            Inches(10.8),
            Inches(0.55),
            desc,
            size=13,
            color=c("muted"),
            font=SANS,
        )
        y += Inches(0.95)


def screenshot_slide(prs, eyebrow, title, caption, image_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, eyebrow, title)
    frame_left = MARGIN_L
    frame_top = Inches(1.55)
    frame_w = Inches(7.4)
    frame_h = Inches(5.2)
    frame = slide.shapes.add_shape(1, frame_left, frame_top, frame_w, frame_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = c("ivory")
    frame.line.color.rgb = c("border")
    frame.line.width = Pt(1.5)

    image_path = SCREENSHOTS / image_name
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), frame_left + Inches(0.08), frame_top + Inches(0.08), width=frame_w - Inches(0.16))
    else:
        add_textbox(
            slide,
            frame_left + Inches(0.3),
            frame_top + Inches(2.2),
            frame_w - Inches(0.6),
            Inches(0.8),
            f"[ Screenshot: {image_name} ]",
            size=14,
            color=c("muted"),
            italic=True,
            font=SANS,
            align=PP_ALIGN.CENTER,
        )

    cap_left = Inches(8.55)
    add_eyebrow(slide, "Product", top=Inches(1.75))
    for i, line in enumerate(caption.split("\n")):
        add_textbox(
            slide,
            cap_left,
            Inches(2.15) + Inches(i * 0.55),
            Inches(4.0),
            Inches(0.5),
            line,
            size=16 if i == 0 else 14,
            italic=(i == 0),
            color=c("forest_deep") if i == 0 else c("muted"),
            font=SERIF if i == 0 else SANS,
        )


def outcomes_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Market", "Potential outcomes")
    scenarios = [
        ("Best case", "Category leader in AI group travel planning — the default way groups plan and book trips."),
        ("Realistic", "[ To be completed ] — meaningful share of US group leisure planning with affiliate + API revenue."),
        ("Worst case", "Niche planning tool for organizers — still saves groups from group-chat purgatory."),
    ]
    x_positions = [MARGIN_L, Inches(4.55), Inches(8.35)]
    for (label, body), x in zip(scenarios, x_positions):
        add_callout_box(slide, x, Inches(1.65), Inches(3.55), Inches(2.2), label, body)


def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, c("forest_deep"))
    add_eyebrow(slide, "Group travel, Handled  ·  Est. 2026", on_dark=True)
    add_textbox(
        slide,
        MARGIN_L,
        Inches(1.85),
        CONTENT_W,
        Inches(1.4),
        "AVANTI",
        size=80,
        color=c("cream"),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(3.35),
        CONTENT_W,
        Inches(0.55),
        "Next-Generation Group Travel",
        size=22,
        color=c("cream"),
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(4.0),
        CONTENT_W,
        Inches(0.45),
        "All the dream. None of the nightmare.",
        size=16,
        color=c("cream"),
        align=PP_ALIGN.CENTER,
        font=SANS,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(5.85),
        CONTENT_W,
        Inches(0.7),
        "Alexandra Goldin, Founder\nInvestor Presentation  ·  Confidential",
        size=13,
        color=c("cream"),
        align=PP_ALIGN.CENTER,
        font=SANS,
    )
    add_footer(slide, on_dark=True)


def team_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Team", "The people")
    add_callout_box(
        slide,
        MARGIN_L,
        Inches(1.65),
        Inches(5.2),
        Inches(1.6),
        "Alexandra Goldin — Founder",
        "UMich '26  ·  linkedin.com/in/alexandragoldin",
    )
    add_callout_box(
        slide,
        Inches(6.35),
        Inches(1.65),
        Inches(5.2),
        Inches(2.4),
        "[ To be completed ]",
        "Founder story\nAdvisors / co-founders\nKey hires planned",
    )


def progress_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = slide_header(slide, "Traction", "Progress to date")
    shipped = [
        "Full planning pipeline live — invite through dining + game time.",
        "Destination matrix, structured voting, and flight analysis shipped.",
        "Booking vault, expense splitting, daily briefings, and essentials live.",
        "Affiliate attribution layer wired (Kayak, Booking, GYG, Expedia/VRBO).",
        "Free beta at avanti.app — try-before-signup at /try.",
    ]
    add_bullets(slide, MARGIN_L, top, Inches(6.2), Inches(4.5), shipped, size=14)
    add_callout_box(
        slide,
        Inches(7.3),
        top,
        Inches(4.25),
        Inches(3.2),
        "[ To be completed ]",
        "Registered users\nTrips created\nDestinations locked\nBooking attribution",
    )


def market_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = slide_header(slide, "Market", "Market size")
    add_bullets(
        slide,
        MARGIN_L,
        top,
        Inches(6.0),
        Inches(3.5),
        [
            "TAM — Global leisure travel: ~$800B–$1T / year.",
            "SAM — US group leisure planning: ~$150–200B addressable.",
            "~40–60M Americans take a group trip annually.",
        ],
        size=15,
    )
    add_callout_box(
        slide,
        Inches(7.3),
        top,
        Inches(4.25),
        Inches(2.0),
        "SOM — Year 3",
        "[ To be completed: GMV, take rate, revenue target ]",
    )


def business_model_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = slide_header(slide, "Business model", "How we make money")
    add_bullets(
        slide,
        MARGIN_L,
        top,
        Inches(6.2),
        Inches(4.0),
        [
            "Today — affiliate commissions on flights, hotels, and activities.",
            "Next — API booking margin via Duffel (flights) and LiteAPI (hotels).",
            "Later — premium concierge tier for high-touch groups.",
            "Free during beta → take-rate at scale.",
        ],
        size=15,
    )
    add_callout_box(
        slide,
        Inches(7.3),
        top,
        Inches(4.25),
        Inches(2.0),
        "Unit economics",
        "[ To be completed: avg trip GMV, take rate, revenue per trip ]",
    )


def ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, c("forest_deep"))
    add_eyebrow(slide, "Fundraising", on_dark=True)
    add_textbox(slide, MARGIN_L, Inches(1.2), CONTENT_W, Inches(0.8), "The ask", size=40, color=c("cream"))
    add_divider(slide, Inches(2.05), color=c("forest"))
    add_textbox(
        slide,
        MARGIN_L,
        Inches(2.8),
        CONTENT_W,
        Inches(1.2),
        "[ To be completed ]",
        size=22,
        color=c("cream"),
        italic=True,
        font=SANS,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(4.2),
        CONTENT_W,
        Inches(1.5),
        "Amount  ·  Use of funds  ·  Runway  ·  Milestones unlocked",
        size=14,
        color=c("cream"),
        font=SANS,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, on_dark=True)


def main():
    SCREENSHOTS.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover_slide(prs)

    headline_slide(
        prs,
        "Problem",
        "Group travel is an absolute nightmare.",
        [
            "Too many people, opinions, and moving parts — and it all lands on one person.",
            "Group chats spiral; spreadsheets multiply; nobody wants to be the bad guy.",
            "Planning stops at ideas — the trip everyone wanted, the planning nobody did.",
        ],
    )

    quote_slide(
        prs,
        "Problem",
        "Tools weren't built for this.",
        '"I\'m honestly up for anything" — but that still means someone has to decide.',
        [
            "Booking apps assume one decider. AI chatbots give lists, not decisions.",
            "Different departure cities, budgets, passports, and date conflicts go untracked.",
            "No structure for votes, feasibility, or getting to a booked trip.",
        ],
    )

    headline_slide(
        prs,
        "Solution",
        "Avanti — the OS for group trips.",
        [
            "AI-native platform built for travel decisions — not another chat thread.",
            "Beachhead: friend groups & families, 4–12 people, US-first.",
            "The only travel app that thinks for everyone — and of everything.",
        ],
        subtitle="More than a chatbot. More than a travel agent.",
    )

    headline_slide(
        prs,
        "Solution",
        "People  ·  Plan  ·  Place",
        [
            "People — bring the group; Avanti keeps the rhythm.",
            "Plan — itineraries, votes, logistics handled before you board.",
            "Place — rooms, tables, and corners worth flying for — booked and confirmed.",
            "Flow: Invite → Brainstorm → Choose → Plan → Game Time.",
        ],
    )

    differentiators_slide(
        prs,
        [
            ("01", "Not another AI chat", "Programmed for travel decisions — 80+ weighted factors per destination."),
            ("02", "Built for groups", "Structured votes, per-person pricing, async windows — nobody is the bad guy."),
            ("03", "Everything in one place", "Research, voting, bookings, vault, briefings — one trip, one home."),
            ("04", "Thinks ahead of you", "Status perks, hidden costs, feasibility gates — before you commit."),
            ("05", "Does the work", "Narrows options, aligns the group, keeps building until wheels up."),
        ],
    )

    screenshot_slide(
        prs,
        "Product",
        "Brainstorm & compare",
        "Destination matrix\nTiered options with honest tradeoffs for your whole group.",
        "brainstorm.png",
    )
    screenshot_slide(
        prs,
        "Product",
        "Decide together",
        "Structured voting\nAsync votes with feasibility gates — so nobody has to be the bad guy.",
        "voting.png",
    )
    screenshot_slide(
        prs,
        "Product",
        "Game time",
        "On-trip companion\nItinerary, bookings vault, briefings, and expense splitting.",
        "gametime.png",
    )

    headline_slide(
        prs,
        "Use cases",
        "Where Avanti wins first",
        [
            "Friend getaways — 6–10 people, mixed budgets, one motivated organizer.",
            "Family reunions — multi-city departures, kids as managed travelers.",
            "Bachelor / milestone trips — date-sensitive, high coordination pain.",
            "Corporate offsites & retreats — B2B expansion path.",
        ],
    )

    headline_slide(
        prs,
        "Benefits",
        "For everyone in the group",
        [
            "Organizer — no more PM, therapist, and bad guy in one.",
            "Members — exactly as much say as you want; personalized cost before commit.",
            "The trip — better than the group chat would find; actually gets booked.",
        ],
    )

    business_model_slide(prs)

    headline_slide(
        prs,
        "Technology",
        "AI-native infrastructure",
        [
            "Decision engine — 80+ destination factors; deal-breakers heavily penalized.",
            "Claude for generation, flight analysis, receipt parsing, inspiration extraction.",
            "Supabase — auth, Postgres, RLS; per-trip affiliate attribution.",
            "Duffel, LiteAPI, GetYourGuide APIs + Kayak / Booking / Expedia affiliates.",
        ],
    )

    headline_slide(
        prs,
        "Go to market",
        "Beachhead & flywheel",
        [
            "Organic — @withavanti on Instagram & TikTok; /try preview funnel.",
            "Word of mouth — every trip invites 4–12 new users.",
            "Campus & network beachhead → creator-led trips → paid + B2B2C.",
            "Flywheel: organizer plans → members join → some become organizers.",
        ],
    )

    market_slide(prs)
    outcomes_slide(prs)

    headline_slide(
        prs,
        "Why now",
        "The moment is right",
        [
            "AI can finally optimize multi-person constraints — impossible with rules engines.",
            "Post-COVID group travel surge — reunions, friend trips, multi-gen travel.",
            "Mobile-native planners expect software, not phone calls to travel agents.",
        ],
    )

    progress_slide(prs)
    team_slide(prs)
    ask_slide(prs)

    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
